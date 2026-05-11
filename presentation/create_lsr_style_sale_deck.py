from pathlib import Path

import requests
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS_DIR = ROOT / "assets"
OUTPUT = ROOT / "lsr_sale_deal_presentation.pptx"

RED = RGBColor(211, 18, 38)
DARK = RGBColor(33, 33, 33)
MID = RGBColor(92, 92, 92)
LIGHT = RGBColor(246, 247, 249)
LINE = RGBColor(223, 226, 231)
WHITE = RGBColor(255, 255, 255)

FONT = "Arial"
FONT_BOLD = "Arial"

PIXABAY_IMAGES = {
    "residential": {
        "url": "https://cdn.pixabay.com/photo/2022/02/06/08/57/residential-complex-6996596_1280.jpg",
        "source": "https://pixabay.com/photos/residential-complex-building-6996596/",
    },
    "signing": {
        "url": "https://cdn.pixabay.com/photo/2021/09/04/21/40/signing-6598540_1280.jpg",
        "source": "https://pixabay.com/photos/signing-paper-document-stamp-deal-6598540/",
    },
    "housing": {
        "url": "https://cdn.pixabay.com/photo/2016/09/05/20/19/housing-1647624_1280.jpg",
        "source": "https://pixabay.com/photos/housing-real-estate-estate-mortgage-1647624/",
    },
}


def ensure_images() -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}
    headers = {"User-Agent": "Mozilla/5.0"}
    for name, meta in PIXABAY_IMAGES.items():
        target = ASSETS_DIR / f"{name}.jpg"
        if not target.exists():
            response = requests.get(meta["url"], headers=headers, timeout=45)
            response.raise_for_status()
            target.write_bytes(response.content)
        paths[name] = target
    return paths


def crop_image(src: Path, name: str, width: int, height: int) -> Path:
    target = ASSETS_DIR / f"{name}_{width}x{height}.jpg"
    image = Image.open(src).convert("RGB")
    cropped = ImageOps.fit(image, (width, height), method=Image.Resampling.LANCZOS)
    cropped.save(target, quality=92)
    return target


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, color, transparency=0, line_color=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = line_color or color
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    line_spacing=None,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BOLD if bold else FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=14, color=DARK, bullet=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_brand(slide, page: str):
    add_rect(slide, Inches(0.55), Inches(0.35), Inches(0.48), Inches(0.48), RED)
    add_text(slide, "ЛСР", Inches(0.605), Inches(0.475), Inches(0.37), Inches(0.16), 8, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "ГРУППА ЛСР", Inches(1.13), Inches(0.42), Inches(1.4), Inches(0.22), 8.5, DARK, True)
    add_text(slide, "практика сделки", Inches(1.13), Inches(0.62), Inches(1.5), Inches(0.18), 7.5, MID)
    add_rect(slide, Inches(11.75), Inches(0.48), Inches(0.78), Inches(0.08), RED)
    add_text(slide, page, Inches(12.58), Inches(0.38), Inches(0.28), Inches(0.22), 8, MID, True, PP_ALIGN.RIGHT)


def add_footer(slide, page: str, photo_source: str):
    add_rect(slide, Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.01), LINE)
    add_text(
        slide,
        f"Фото: Pixabay ({photo_source})",
        Inches(0.55),
        Inches(7.12),
        Inches(6.2),
        Inches(0.18),
        6.5,
        MID,
    )
    add_text(slide, page, Inches(12.2), Inches(7.12), Inches(0.6), Inches(0.18), 6.5, MID, True, PP_ALIGN.RIGHT)


def add_label(slide, text, x, y, w):
    add_rect(slide, x, y, w, Inches(0.32), RED, radius=False)
    add_text(slide, text, x + Inches(0.13), y + Inches(0.075), w - Inches(0.26), Inches(0.12), 7.5, WHITE, True)


def add_card(slide, number, title, body, x, y, w, h):
    add_rect(slide, x, y, w, h, WHITE, line_color=LINE, radius=True)
    add_text(slide, number, x + Inches(0.22), y + Inches(0.2), Inches(0.5), Inches(0.2), 11, RED, True)
    add_text(slide, title, x + Inches(0.22), y + Inches(0.55), w - Inches(0.44), Inches(0.3), 15, DARK, True)
    add_text(slide, body, x + Inches(0.22), y + Inches(0.98), w - Inches(0.44), h - Inches(1.1), 10.8, MID, line_spacing=1.1)


def title_slide(prs: Presentation, images: dict[str, Path]):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_brand(slide, "01")

    photo = crop_image(images["residential"], "slide1_photo", 930, 1080)
    slide.shapes.add_picture(str(photo), Inches(7.15), Inches(0.9), Inches(5.55), Inches(5.92))
    add_rect(slide, Inches(7.15), Inches(0.9), Inches(5.55), Inches(5.92), RED, transparency=78)
    add_rect(slide, Inches(7.0), Inches(1.18), Inches(0.16), Inches(5.08), RED)

    add_label(slide, "ПОКУПКА / ПРОДАЖА", Inches(0.72), Inches(1.15), Inches(1.85))
    add_text(
        slide,
        "Как заключить сделку купли-продажи недвижимости",
        Inches(0.7),
        Inches(1.65),
        Inches(6.05),
        Inches(1.6),
        31,
        DARK,
        True,
        line_spacing=0.9,
    )
    add_text(
        slide,
        "Короткий маршрут: от проверки объекта и согласования условий до регистрации права и передачи ключей.",
        Inches(0.72),
        Inches(3.42),
        Inches(5.75),
        Inches(0.7),
        14,
        MID,
        line_spacing=1.1,
    )

    chips = [
        ("01", "Проверить объект и участников"),
        ("02", "Зафиксировать условия договора"),
        ("03", "Зарегистрировать переход права"),
    ]
    y = Inches(4.45)
    for num, text in chips:
        add_rect(slide, Inches(0.72), y, Inches(0.43), Inches(0.43), RED)
        add_text(slide, num, Inches(0.82), y + Inches(0.12), Inches(0.24), Inches(0.14), 7.5, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, text, Inches(1.32), y + Inches(0.09), Inches(4.8), Inches(0.18), 12.5, DARK, True)
        y += Inches(0.6)

    add_rect(slide, Inches(8.0), Inches(5.55), Inches(3.9), Inches(0.88), DARK, transparency=10)
    add_text(
        slide,
        "Принцип: деньги, документы и сроки движутся по одному согласованному сценарию.",
        Inches(8.28),
        Inches(5.82),
        Inches(3.35),
        Inches(0.28),
        11,
        WHITE,
        True,
        line_spacing=1.05,
    )
    add_footer(slide, "01 / 03", PIXABAY_IMAGES["residential"]["source"])


def prep_slide(prs: Presentation, images: dict[str, Path]):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_brand(slide, "02")
    add_label(slide, "ДО ДОГОВОРА", Inches(0.72), Inches(1.02), Inches(1.55))
    add_text(slide, "Что нужно подготовить", Inches(0.72), Inches(1.48), Inches(5.3), Inches(0.55), 30, DARK, True)
    add_text(
        slide,
        "Главная задача подготовки — заранее убрать юридические, финансовые и организационные неопределенности.",
        Inches(0.72),
        Inches(2.1),
        Inches(6.1),
        Inches(0.48),
        12.5,
        MID,
        line_spacing=1.1,
    )

    photo = crop_image(images["signing"], "slide2_photo", 960, 720)
    slide.shapes.add_picture(str(photo), Inches(7.45), Inches(1.08), Inches(4.75), Inches(3.56))
    add_rect(slide, Inches(7.2), Inches(0.9), Inches(0.12), Inches(3.95), RED)
    add_rect(slide, Inches(10.85), Inches(4.32), Inches(1.1), Inches(0.38), RED)
    add_text(slide, "проверка", Inches(10.98), Inches(4.42), Inches(0.8), Inches(0.12), 7, WHITE, True, PP_ALIGN.CENTER)

    add_card(
        slide,
        "01",
        "Объект",
        "Выписка ЕГРН, обременения, история владения, перепланировки, задолженности и основания права.",
        Inches(0.72),
        Inches(3.05),
        Inches(3.85),
        Inches(1.95),
    )
    add_card(
        slide,
        "02",
        "Стороны",
        "Паспорта, семейный статус, согласия, полномочия представителя, дееспособность и действительность доверенностей.",
        Inches(4.78),
        Inches(3.05),
        Inches(3.85),
        Inches(1.95),
    )
    add_card(
        slide,
        "03",
        "Расчеты",
        "Цена, аванс или задаток, ипотека, аккредитив, эскроу или ячейка, условия раскрытия денег.",
        Inches(0.72),
        Inches(5.22),
        Inches(3.85),
        Inches(1.33),
    )
    add_card(
        slide,
        "04",
        "Сроки",
        "Дата подписания, подача на регистрацию, освобождение объекта, передача ключей и ответственность за срыв.",
        Inches(4.78),
        Inches(5.22),
        Inches(3.85),
        Inches(1.33),
    )
    add_footer(slide, "02 / 03", PIXABAY_IMAGES["signing"]["source"])


def closing_slide(prs: Presentation, images: dict[str, Path]):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE)
    add_brand(slide, "03")

    photo = crop_image(images["housing"], "slide3_photo", 860, 1080)
    slide.shapes.add_picture(str(photo), Inches(0.0), Inches(0.0), Inches(4.65), Inches(7.5))
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(4.65), Inches(7.5), DARK, transparency=54)
    add_rect(slide, Inches(3.98), Inches(0.0), Inches(0.24), Inches(7.5), RED)

    add_text(slide, "Финальный этап", Inches(0.62), Inches(1.12), Inches(2.9), Inches(0.26), 12, WHITE, True)
    add_text(slide, "Подписание,\nрасчеты и регистрация", Inches(0.62), Inches(1.62), Inches(3.15), Inches(1.2), 25, WHITE, True, line_spacing=0.9)
    add_text(
        slide,
        "Право собственности возникает после государственной регистрации перехода права.",
        Inches(0.62),
        Inches(4.95),
        Inches(2.9),
        Inches(0.58),
        11,
        WHITE,
        line_spacing=1.08,
    )

    add_label(slide, "ПОРЯДОК ДЕЙСТВИЙ", Inches(5.05), Inches(1.02), Inches(1.95))
    add_text(slide, "Четыре шага к закрытой сделке", Inches(5.05), Inches(1.48), Inches(6.3), Inches(0.52), 29, DARK, True)

    steps = [
        ("1", "Подписать договор", "Существенные условия, цена, порядок оплаты, сроки освобождения и передачи объекта."),
        ("2", "Подать документы", "МФЦ, Росреестр или электронная регистрация через банк/нотариуса."),
        ("3", "Раскрыть расчеты", "Деньги передаются по условиям аккредитива, эскроу, ячейки или иной безопасной схемы."),
        ("4", "Принять объект", "Выписка ЕГРН, акт приема-передачи, ключи, счетчики и подтверждение взаиморасчетов."),
    ]
    y = Inches(2.35)
    for num, title, body in steps:
        add_rect(slide, Inches(5.05), y, Inches(0.48), Inches(0.48), RED)
        add_text(slide, num, Inches(5.2), y + Inches(0.13), Inches(0.18), Inches(0.14), 8.5, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, title, Inches(5.76), y + Inches(0.02), Inches(3.2), Inches(0.22), 13.5, DARK, True)
        add_text(slide, body, Inches(5.76), y + Inches(0.34), Inches(5.8), Inches(0.35), 9.3, MID, line_spacing=1.05)
        y += Inches(0.86)

    add_rect(slide, Inches(9.05), Inches(4.95), Inches(3.0), Inches(1.45), RED, radius=True)
    add_text(slide, "Финальный чек-лист", Inches(9.32), Inches(5.16), Inches(2.45), Inches(0.2), 12.5, WHITE, True)
    add_multiline(
        slide,
        [
            "договор и акт подписаны",
            "расчеты подтверждены",
            "выписка ЕГРН получена",
        ],
        Inches(9.32),
        Inches(5.52),
        Inches(2.45),
        Inches(0.62),
        size=8.5,
        color=WHITE,
        bullet=True,
    )
    add_text(
        slide,
        "Материал носит информационный характер; условия сделки нужно сверять с юристом, банком или нотариусом.",
        Inches(5.05),
        Inches(6.7),
        Inches(5.95),
        Inches(0.24),
        7.2,
        MID,
    )
    add_footer(slide, "03 / 03", PIXABAY_IMAGES["housing"]["source"])


def build_deck():
    images = ensure_images()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs, images)
    prep_slide(prs, images)
    closing_slide(prs, images)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
